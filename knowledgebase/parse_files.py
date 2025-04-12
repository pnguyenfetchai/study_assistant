import fitz  # PyMuPDF for PDFs
import docx
import pytesseract  # OCR for images
import pandas as pd  # For Excel files
import os
from PIL import Image
from pptx import Presentation  # For PowerPoint files

def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF file."""
    text = ""
    try:
        doc = fitz.open(pdf_path)
        for page in doc:
            text += page.get_text("text") + "\n"
    except Exception as e:
        print(f"⚠️ Error processing PDF {pdf_path}: {e}")
    return text


def extract_text_from_docx(docx_path):
    """Extract text from a Word document."""
    doc = docx.Document(docx_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_txt(txt_path):
    """Extract text from a TXT file."""
    with open(txt_path, "r", encoding="utf-8") as file:
        return file.read()

def extract_text_from_pptx(pptx_path):
    """Extract text from a PowerPoint (.pptx) file."""
    presentation = Presentation(pptx_path)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_images(image_path):
    """Extract text from an image using Tesseract OCR."""
    try:
        image = Image.open(image_path)
        return pytesseract.image_to_string(image)
    except Exception as e:
        return f"Error processing image {image_path}: {str(e)}"

def extract_text_from_excel(excel_path):
    """Extract text from an Excel file (.xlsx, .xls)."""
    try:
        df = pd.read_excel(excel_path, engine='openpyxl')
        return df.to_string(index=False)  # Convert dataframe to string
    except Exception as e:
        return f"Error processing Excel file {excel_path}: {str(e)}"

def extract_text_from_files(directory):
    """Extract text from all supported file types in a directory."""
    docs = []
    for root, _, files in os.walk(directory):
        for file in files:
            file_path = os.path.join(root, file)
            text = None  # Initialize empty text content
            
            if file.endswith(".pdf"):
                text = extract_text_from_pdf(file_path)
            elif file.endswith(".docx"):
                text = extract_text_from_docx(file_path)
            elif file.endswith(".txt"):
                text = extract_text_from_txt(file_path)
            elif file.endswith(".pptx"):
                text = extract_text_from_pptx(file_path)
            elif file.endswith((".png", ".jpg", ".jpeg")):
                text = extract_text_from_images(file_path)
            elif file.endswith((".xlsx", ".xls")):
                text = extract_text_from_excel(file_path)
            else:
                continue  # Skip unsupported files
            
            if text:  
                relative_path = os.path.relpath(file_path, directory)
                course_folder = relative_path.split(os.sep)[0]
                docs.append(f"Course: {course_folder}, File: {file}, Content: {text}")

    print("dep trai qua di: ", docs)
    return docs

if __name__ == "__main__":
    docs = extract_text_from_files("course_files")
    for doc in docs:
        print("📄", doc[:300])
